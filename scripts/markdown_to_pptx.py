"""
Simple Markdown -> PowerPoint converter tailored to the project's `docs/papers_summaries.md`.

This script splits the Markdown into slides by top-level headings and basic separators,
adds title and bullet slides, and embeds images referenced in the markdown (relative paths).

Usage:
    python3 scripts/markdown_to_pptx.py docs/papers_summaries.md docs/papers_summaries.pptx

This script has minimal dependencies: python-pptx and markdown (optional). It does not
attempt to fully render complex Markdown (tables, math, etc.) but handles headings, paragraphs,
lists, and image links.
"""

import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def parse_markdown(md_text):
    """Parse markdown into a list of slides. Each slide is a dict with 'title' and 'content' list."""
    lines = md_text.splitlines()
    slides = []
    cur = {'title': None, 'content': []}

    def flush():
        if cur['title'] or cur['content']:
            slides.append({'title': cur['title'], 'content': cur['content'][:]})
            cur['title'] = None
            cur['content'].clear()

    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        # Title slide marker - first level heading
        m = re.match(r"^#{1} \s*(.+)", line)
        if m:
            # start a new slide
            flush()
            cur['title'] = m.group(1).strip()
            i += 1
            continue
        # Secondary headings become section headers or slide titles
        m2 = re.match(r"^#{2} \s*(.+)", line)
        if m2:
            # start a new slide with this as title
            flush()
            cur['title'] = m2.group(1).strip()
            i += 1
            continue
        # image line ![](path) or ![alt](path)
        mimg = re.search(r"!\[.*?\]\(([^)]+)\)", line)
        if mimg:
            cur['content'].append({'type': 'image', 'path': mimg.group(1).strip()})
            i += 1
            continue
        # bullet list
        mb = re.match(r"^\s*[-*+]\s+(.+)", line)
        if mb:
            cur['content'].append({'type': 'bullet', 'text': mb.group(1).strip()})
            i += 1
            continue
        # table line or separator - skip elaborate parsing, treat as paragraph
        if line.strip().startswith('|') or line.strip().startswith('---'):
            # include as paragraph
            cur['content'].append({'type': 'para', 'text': line.strip()})
            i += 1
            continue
        # normal paragraph
        if line.strip():
            cur['content'].append({'type': 'para', 'text': line.strip()})
        else:
            # blank line might indicate slide break if current has content
            # but we'll just keep accumulating until next heading
            pass
        i += 1

    flush()
    return slides


def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    if subtitle:
        try:
            subtitle_tf = slide.placeholders[1]
            subtitle_tf.text = subtitle
        except Exception:
            pass
    return slide


def add_section_slide(prs, title, contents, assets_dir):
    # use Title and Content layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title if title else ''
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for item in contents:
        if item['type'] == 'bullet':
            p = body.add_paragraph()
            p.text = item['text']
            p.level = 0
            p.font.size = Pt(18)
        elif item['type'] == 'para':
            p = body.add_paragraph()
            p.text = item['text']
            p.level = 0
            p.font.size = Pt(14)
        elif item['type'] == 'image':
            img_path = Path(item['path'])
            if not img_path.is_absolute():
                img_path = assets_dir / img_path
            if img_path.exists():
                # add image below the text box
                left = Inches(1)
                top = Inches(2.2)
                width = Inches(6)
                slide.shapes.add_picture(str(img_path), left, top, width=width)
            else:
                p = body.add_paragraph()
                p.text = f"[Missing image: {item['path']}]"
    return slide


def convert(md_path: Path, pptx_path: Path):
    md_text = md_path.read_text(encoding='utf-8')
    slides = parse_markdown(md_text)
    prs = Presentation()
    assets_dir = md_path.parent

    if slides:
        # If first slide is a title, make a title slide
        first = slides[0]
        if first['title']:
            # Use first title as title slide and remove it from list
            add_title_slide(prs, first['title'])
            # also include its content as a following slide if present
            if first['content']:
                add_section_slide(prs, first['title'], first['content'], assets_dir)
            slides = slides[1:]

    for s in slides:
        title = s['title'] or ''
        add_section_slide(prs, title, s['content'], assets_dir)

    prs.save(str(pptx_path))


def main():
    if len(sys.argv) < 3:
        print('Usage: python3 scripts/markdown_to_pptx.py INPUT.md OUTPUT.pptx')
        sys.exit(1)
    md_path = Path(sys.argv[1])
    pptx_path = Path(sys.argv[2])
    if not md_path.exists():
        print('Input markdown not found:', md_path)
        sys.exit(1)
    convert(md_path, pptx_path)
    print('Saved PPTX to', pptx_path)


if __name__ == '__main__':
    main()
